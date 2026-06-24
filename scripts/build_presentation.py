"""Build the PowerPoint presentation from PRESENTATION.md."""

from __future__ import annotations

import argparse
from pathlib import Path


def _parse_markdown_slides(markdown_text: str) -> list[tuple[str, list[str]]]:
    slides: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_bullets: list[str] = []

    for raw_line in markdown_text.splitlines():
        line = raw_line.strip()
        if line.startswith('## '):
            if current_title is not None:
                slides.append((current_title, current_bullets))
            current_title = line.removeprefix('## ').strip()
            if ': ' in current_title:
                current_title = current_title.split(': ', 1)[1]
            current_bullets = []
        elif line.startswith('- ') and current_title is not None:
            current_bullets.append(line.removeprefix('- ').strip())

    if current_title is not None:
        slides.append((current_title, current_bullets))

    return slides


def _first_markdown_title(markdown_text: str) -> str:
    for raw_line in markdown_text.splitlines():
        line = raw_line.strip()
        if line.startswith('# '):
            return line.removeprefix('# ').strip()
    return 'Генериране на описание за изображение'


def _add_text(slide, left, top, width, height, text, size, *, bold=False, italic=False, align=None):
    from pptx.util import Pt

    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align is not None:
        paragraph.alignment = align

    for run in paragraph.runs:
        run.font.name = 'Calibri'
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
    return shape


def _add_footer(slide, title: str, index: int, total: int) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    _add_text(
        slide,
        Inches(0.55),
        Inches(7.08),
        Inches(8.5),
        Inches(0.25),
        title,
        8,
        italic=True,
        align=PP_ALIGN.LEFT,
    )
    _add_text(
        slide,
        Inches(11.8),
        Inches(7.08),
        Inches(0.9),
        Inches(0.25),
        f'{index}/{total}',
        8,
        align=PP_ALIGN.RIGHT,
    )


def _add_title_slide(presentation, project_title: str, first_slide_bullets: list[str]) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    _add_text(
        slide,
        Inches(0.7),
        Inches(0.35),
        Inches(12.0),
        Inches(0.45),
        'Софийски университет „Св. Климент Охридски“',
        21,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        Inches(1.15),
        Inches(1.65),
        Inches(11.1),
        Inches(1.0),
        project_title,
        36,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        Inches(3.8),
        Inches(3.0),
        Inches(5.8),
        Inches(0.55),
        'Курсов проект',
        30,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    footer_lines = [
        'Факултет по математика и информатика',
        'Студент: Никола',
        'Факултетен номер: [попълва се]',
        'Учебен план: Изкуствен интелект',
        'Дисциплина: Дълбоко обучение',
    ]
    if first_slide_bullets:
        footer_lines.append(first_slide_bullets[-1])

    for offset, line in enumerate(footer_lines):
        _add_text(
            slide,
            Inches(2.6),
            Inches(4.6 + offset * 0.33),
            Inches(8.2),
            Inches(0.28),
            line,
            15,
            italic=True,
            align=PP_ALIGN.CENTER,
        )


def _add_content_slide(presentation, project_title: str, slide_title: str, bullets: list[str], index: int, total: int) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    _add_text(
        slide,
        Inches(0.65),
        Inches(0.38),
        Inches(12.1),
        Inches(0.5),
        slide_title,
        25,
        bold=True,
    )

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.65),
        Inches(1.02),
        Inches(12.0),
        Inches(0.02),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0, 0, 0)
    line.line.color.rgb = RGBColor(0, 0, 0)

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.6), Inches(5.55))
    frame = body.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.05)
    frame.margin_right = Inches(0.05)

    bullet_count = max(len(bullets), 1)
    font_size = 20 if bullet_count <= 5 else 18
    space_after = 10 if bullet_count <= 5 else 6

    for bullet_index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(space_after)
        paragraph.font.name = 'Calibri'
        paragraph.font.size = Pt(font_size)

    _add_footer(slide, project_title, index, total)


def build_presentation(input_path: Path, output_path: Path) -> None:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise SystemExit(
            'Missing dependency: python-pptx. Install dev requirements with '
            '`python -m pip install -r requirements-dev.txt`.'
        ) from exc

    markdown_text = input_path.read_text(encoding='utf-8')
    project_title = _first_markdown_title(markdown_text)
    slides = _parse_markdown_slides(markdown_text)
    if not slides:
        raise ValueError(f'No slides found in {input_path}')

    presentation = Presentation()
    presentation.slide_width = 12192000
    presentation.slide_height = 6858000

    _add_title_slide(presentation, project_title, slides[0][1])

    content_slides = slides[1:]
    total = len(content_slides) + 1
    for slide_number, (title, bullets) in enumerate(content_slides, start=2):
        _add_content_slide(presentation, project_title, title, bullets, slide_number, total)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description='Build PowerPoint slides from PRESENTATION.md.')
    parser.add_argument('--input', type=Path, default=Path('PRESENTATION.md'))
    parser.add_argument('--output', type=Path, default=Path('presentation/image_captioning_presentation.pptx'))
    args = parser.parse_args()

    build_presentation(args.input, args.output)
    print(f'Saved presentation to {args.output}')


if __name__ == '__main__':
    main()
